#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Tool Calls 流程分析 - 专业商业演示文稿生成器 v2.0
更加商业化、专业化的设计
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsmap
import os

class ProfessionalPPTGenerator:
    """生成专业商业级别的PPT"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # 高级商业配色方案 (参考McKinsey/BCG风格)
        self.colors = {
            'primary': RGBColor(20, 50, 100),          # 深海蓝
            'secondary': RGBColor(0, 120, 215),         # 专业蓝
            'accent': RGBColor(255, 140, 0),            # 活力橙
            'success': RGBColor(32, 156, 93),           # 商务绿
            'warning': RGBColor(255, 192, 0),           # 警示黄
            'danger': RGBColor(220, 53, 69),            # 错误红
            'dark': RGBColor(33, 37, 41),               # 深黑
            'gray': RGBColor(108, 117, 125),            # 灰色
            'light_gray': RGBColor(233, 236, 239),      # 浅灰
            'white': RGBColor(255, 255, 255),
            'ice_blue': RGBColor(230, 242, 255),        # 冰蓝
            'navy': RGBColor(0, 32, 96),                # 海军蓝
            'teal': RGBColor(0, 128, 128),              # 青色
        }

    def add_gradient_background(self, slide, color1, color2):
        """添加渐变背景"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color1
        shape.line.fill.background()
        shape.z_order = 0
        return shape

    def add_slide_number(self, slide, number):
        """添加页码"""
        box = slide.shapes.add_textbox(
            Inches(12.5), Inches(6.8),
            Inches(0.6), Inches(0.4)
        )
        frame = box.text_frame
        frame.text = str(number)
        para = frame.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        para.font.size = Pt(10)
        para.font.color.rgb = self.colors['gray']

    def add_header_bar(self, slide, text=""):
        """添加顶部装饰条"""
        # 主条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, self.prs.slide_width, Inches(0.12)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['primary']
        bar.line.fill.background()

        # 副条
        bar2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(0.12), self.prs.slide_width, Inches(0.03)
        )
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = self.colors['accent']
        bar2.line.fill.background()

    def create_card(self, slide, left, top, width, height, title, content,
                    bg_color=None, border_color=None, icon=""):
        """创建专业卡片样式"""
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        if bg_color:
            card.fill.solid()
            card.fill.fore_color.rgb = bg_color
        else:
            card.fill.solid()
            card.fill.fore_color.rgb = self.colors['white']

        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(2)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.12),
            Inches(width - 0.3), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{icon} {title}" if icon else title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.55),
            Inches(width - 0.3), Inches(height - 0.65)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = content
        content_para = content_frame.paragraphs[0]
        content_para.font.size = Pt(12)
        content_para.font.color.rgb = self.colors['dark']
        content_para.line_spacing = 1.4

        return card

    def create_numbered_card(self, slide, number, left, top, width, height,
                           title, content, color):
        """创建带编号的卡片"""
        # 背景卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors['white']
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # 编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + 0.1), Inches(top + 0.1),
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # 编号文字
        num_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.18),
            Inches(0.5), Inches(0.35)
        )
        num_frame = num_box.text_frame
        num_frame.text = str(number)
        num_para = num_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(20)
        num_para.font.bold = True
        num_para.font.color.rgb = self.colors['white']

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.7), Inches(top + 0.15),
            Inches(width - 0.85), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = color

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.6),
            Inches(width - 0.3), Inches(height - 0.7)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = content
        content_para = content_frame.paragraphs[0]
        content_para.font.size = Pt(11)
        content_para.font.color.rgb = self.colors['dark']
        content_para.line_spacing = 1.3

        return card

    def add_title_slide(self):
        """封面页 - 专业设计"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 渐变背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['navy']
        bg.line.fill.background()

        # 装饰圆圈
        for i in range(5):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(12 - i * 2.5), Inches(6 - i * 1.2),
                Inches(1.5 + i * 0.5), Inches(1.5 + i * 0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors['accent']
            circle.fill.fore_color.brightness = 0.3 - i * 0.05
            circle.line.fill.background()

        # 主标题框
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.2), Inches(11.333), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        # 主标题
        p = title_frame.paragraphs[0]
        p.text = "Tool Calls"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(64)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.space_after = Pt(10)

        # 副标题
        p = title_frame.add_paragraph()
        p.text = "流程分析报告"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.colors['accent']

        # 描述
        desc_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5), Inches(11.333), Inches(1)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = "完整数据流与实现机制解析 | Technical Documentation"
        desc_para = desc_frame.paragraphs[0]
        desc_para.alignment = PP_ALIGN.CENTER
        desc_para.font.size = Pt(18)
        desc_para.font.color.rgb = RGBColor(180, 200, 230)

        # 底部装饰线
        line1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.5), Inches(5.8), Inches(4.333), Inches(0.08)
        )
        line1.fill.solid()
        line1.fill.fore_color.rgb = self.colors['accent']
        line1.line.fill.background()

        # 日期标签
        date_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(6.5), Inches(1.5), Inches(0.4)
        )
        date_frame = date_box.text_frame
        date_frame.text = "2025"
        date_para = date_frame.paragraphs[0]
        date_para.alignment = PP_ALIGN.RIGHT
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = self.colors['gray']

        self.add_slide_number(slide, "01")

    def add_agenda_slide(self):
        """目录页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self.add_header_bar(slide)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "目录 | Contents"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        agendas = [
            ("01", "数据流概览", "Data Flow Overview"),
            ("02", "数据获取机制", "Data Acquisition"),
            ("03", "工具调用提取", "Tool Calls Extraction"),
            ("04", "用户交互流程", "User Interaction"),
            ("05", "问题诊断方案", "Troubleshooting"),
            ("06", "关键代码索引", "Code Reference"),
            ("07", "总结与建议", "Summary & Recommendations"),
        ]

        y_start = 1.8
        item_height = 0.65

        for i, (num, title_cn, title_en) in enumerate(agendas):
            y_pos = y_start + i * item_height

            # 编号
            num_box = slide.shapes.add_textbox(
                Inches(1), Inches(y_pos), Inches(0.8), Inches(0.5)
            )
            num_frame = num_box.text_frame
            num_frame.text = num
            num_para = num_frame.paragraphs[0]
            num_para.font.size = Pt(24)
            num_para.font.bold = True
            num_para.font.color.rgb = self.colors['accent']

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(2), Inches(y_pos + 0.05), Inches(5), Inches(0.45)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"{title_cn}"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(20)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['primary']

            # 英文标题
            en_box = slide.shapes.add_textbox(
                Inches(2), Inches(y_pos + 0.35), Inches(8), Inches(0.25)
            )
            en_frame = en_box.text_frame
            en_frame.text = title_en
            en_para = en_frame.paragraphs[0]
            en_para.font.size = Pt(12)
            en_para.font.color.rgb = self.colors['gray']
            en_para.font.italic = True

            # 装饰点
            if i < len(agendas) - 1:
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(1.3), Inches(y_pos + 0.5), Inches(0.05), Inches(0.05)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = self.colors['gray']
                dot.line.fill.background()

        # 右侧装饰
        for i in range(3):
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(11.5), Inches(2 + i * 1.2),
                Inches(0.08), Inches(0.8)
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = self.colors['secondary']
            rect.line.fill.background()

        self.add_slide_number(slide, "02")

    def add_data_flow_slide(self):
        """数据流概览"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self.add_header_bar(slide)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "01 数据流概览 | Data Flow Overview"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.95), Inches(12.333), Inches(0.4)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "完整的端到端数据流向分析"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(14)
        subtitle_para.font.color.rgb = self.colors['gray']

        # 流程节点
        nodes = [
            ("后端 API", "Backend API", 0.5, 1.8, self.colors['primary']),
            ("GET /api/v1/session/\n{sessionId}/conversation", "RESTful API", 2.5, 1.8, self.colors['navy']),
            ("getConversation\nHistory()", "Data Fetch", 5.5, 1.8, self.colors['secondary']),
            ("数据转换\n& 格式化", "Transform", 8.5, 1.8, self.colors['accent']),
            ("ChatInterface\n组件", "UI Layer", 0.5, 3.8, self.colors['primary']),
            ("useEffect\n监听", "State Mgmt", 3.5, 3.8, self.colors['navy']),
            ("工具调用\n提取", "Extraction", 6.5, 3.8, self.colors['secondary']),
            ("UI 渲染\nTool Calls", "Display", 9.5, 3.8, self.colors['accent']),
        ]

        for title, subtitle, x, y, color in nodes:
            # 节点框
            node = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(2.5), Inches(1.5)
            )
            node.fill.solid()
            node.fill.fore_color.rgb = color
            node.line.fill.background()

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 0.3),
                Inches(2.2), Inches(0.6)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.size = Pt(13)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['white']

            # 副标题
            sub_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 1),
                Inches(2.2), Inches(0.35)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.alignment = PP_ALIGN.CENTER
            sub_para.font.size = Pt(10)
            sub_para.font.color.rgb = RGBColor(220, 230, 255)

        # 连接箭头 (第一行)
        for i in range(3):
            x = 3 + i * 3
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x), Inches(2.35),
                Inches(0.3), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = self.colors['gray']
            arrow.line.fill.background()

        # 连接箭头 (第二行)
        for i in range(3):
            x = 3 + i * 3
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x), Inches(4.35),
                Inches(0.3), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = self.colors['gray']
            arrow.line.fill.background()

        # 垂直连接
        down_arrow1 = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(1.5), Inches(3.3),
            Inches(0.2), Inches(0.3)
        )
        down_arrow1.fill.solid()
        down_arrow1.fill.fore_color.rgb = self.colors['gray']
        down_arrow1.line.fill.background()

        down_arrow2 = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(6.5), Inches(3.3),
            Inches(0.2), Inches(0.3)
        )
        down_arrow2.fill.solid()
        down_arrow2.fill.fore_color.rgb = self.colors['gray']
        down_arrow2.line.fill.background()

        # 底部说明
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2)
        )
        info_frame = info_box.text_frame
        info_frame.word_wrap = True
        info_frame.text = "📊 关键特性：数据完整性保障 • 双向数据流 • 状态管理 • 实时更新"
        info_para = info_frame.paragraphs[0]
        info_para.alignment = PP_ALIGN.CENTER
        info_para.font.size = Pt(14)
        info_para.font.color.rgb = self.colors['dark']
        info_para.font.bold = True

        self.add_slide_number(slide, "03")

    def add_data_acquisition_slide(self):
        """数据获取机制"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self.add_header_bar(slide)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(6), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "02 数据获取机制 | Data Acquisition"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 左侧 - API信息
        left_cards = [
            ("🔌 API 端点", "GET /api/v1/session/{sessionId}/conversation"),
            ("📂 代码位置", "agentService.ts:390-458"),
            ("⚡ 触发时机", "• 初始化时\n• 切换会话时\n• 外部会话切换"),
        ]

        y_start = 1.5
        for i, (title, content) in enumerate(left_cards):
            y_pos = y_start + i * 1.35
            self.create_card(
                slide, 0.5, y_pos, 5.8, 1.2,
                title, content,
                self.colors['ice_blue'],
                self.colors['secondary']
            )

        # 右侧 - 数据结构
        right_title = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.5), Inches(6), Inches(0.4)
        )
        right_title_frame = right_title.text_frame
        right_title_frame.text = "📦 返回数据结构"
        p = right_title_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # 代码框
        code_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), Inches(1.95),
            Inches(6), Inches(2.7)
        )
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = self.colors['dark']
        code_box.line.color.rgb = self.colors['secondary']
        code_box.line.width = Pt(2)

        code_frame = code_box.text_frame
        code_frame.word_wrap = True
        code_frame.margin_top = Inches(0.15)
        code_frame.margin_left = Inches(0.15)
        code_frame.margin_right = Inches(0.15)

        code_text = '''{
  "session_id": "xxx",
  "messages": [{
    "id": "ai-3",
    "conversation_turn_id": "98a93fe3...",
    "tool_calls": [{
      "tool_use_id": "call_xxx",
      "tool_name": "WebSearch",
      "tool_input": {...},
      "conversation_turn_id": "98a93fe3..."
    }]
  }]
}'''

        p = code_frame.paragraphs[0]
        p.text = code_text
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 255, 100)
        p.font.name = 'Courier New'

        # 关键字段标注
        key_fields = slide.shapes.add_textbox(
            Inches(6.8), Inches(4.8), Inches(6), Inches(1.2)
        )
        key_frame = key_fields.text_frame
        key_frame.word_wrap = True

        fields = [
            ("✅ conversation_turn_id", "对话轮次标识符，用于关联消息"),
            ("✅ tool_calls 数组", "包含该消息的所有工具调用"),
            ("✅ 双重ID保障", "消息级和工具级都有turn_id")
        ]

        for i, (field, desc) in enumerate(fields):
            p = key_frame.add_paragraph() if i > 0 else key_frame.paragraphs[0]
            p.text = f"{field} - {desc}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['success']
            p.font.bold = True
            p.space_before = Pt(4) if i > 0 else Pt(0)

        # 底部转换逻辑
        transform_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.2),
            Inches(12.333), Inches(1)
        )
        transform_box.fill.solid()
        transform_box.fill.fore_color.rgb = self.colors['light_gray']
        transform_box.line.color.rgb = self.colors['accent']
        transform_box.line.width = Pt(2)

        transform_title = slide.shapes.add_textbox(
            Inches(0.7), Inches(6.35), Inches(2), Inches(0.3)
        )
        tt_frame = transform_title.text_frame
        tt_frame.text = "🔄 数据转换"
        p = tt_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['accent']

        transform_content = slide.shapes.add_textbox(
            Inches(2.8), Inches(6.35), Inches(9.8), Inches(0.7)
        )
        tc_frame = transform_content.text_frame
        tc_frame.word_wrap = True
        tc_frame.text = "保留 conversation_turn_id 和 tool_calls 字段，确保数据完整性不丢失"
        p = tc_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['dark']

        self.add_slide_number(slide, "04")

    def add_extraction_slide(self):
        """工具调用提取流程"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self.add_header_bar(slide)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "03 工具调用提取 | Tool Calls Extraction"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 三步流程
        steps = [
            (1, "确定目标轮次",
             "检查 selectedTurnId\n• 有值：使用用户选择\n• 无值：使用最新AI消息",
             self.colors['secondary']),
            (2, "遍历匹配",
             "遍历所有消息的 tool_calls\n• 匹配 conversation_turn_id\n• 字符串比较确保一致性",
             self.colors['accent']),
            (3, "更新状态",
             "setToolCalls(extracted)\n• 触发UI重新渲染\n• 显示选中轮次工具",
             self.colors['success']),
        ]

        y_start = 1.6
        box_width = 3.8
        box_height = 2.2

        for num, title, content, color in steps:
            x = 0.5 + (num - 1) * 4.2
            self.create_numbered_card(
                slide, num, x, y_start, box_width, box_height,
                title, content, color
            )

        # 核心代码
        core_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.4)
        )
        ct_frame = core_title.text_frame
        ct_frame.text = "⚙️ 核心匹配逻辑 (ChatInterface.tsx:1368-1534)"
        p = ct_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        core_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(4.7),
            Inches(12.333), Inches(1.5)
        )
        core_box.fill.solid()
        core_box.fill.fore_color.rgb = self.colors['dark']
        core_box.line.color.rgb = self.colors['secondary']
        core_box.line.width = Pt(2)

        core_frame = core_box.text_frame
        core_frame.word_wrap = True
        core_frame.margin_left = Inches(0.3)

        core_code = '''const toolTurnIdStr = String(toolCallTurnId || '');
const targetTurnIdStr = String(targetTurnId || '');
if (toolTurnIdStr === targetTurnIdStr) {
    extractedToolCalls.push({...});
}
setToolCalls(extractedToolCalls);'''

        p = core_frame.paragraphs[0]
        p.text = core_code
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(100, 255, 100)
        p.font.name = 'Courier New'

        # 依赖说明
        dep_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.4),
            Inches(5.5), Inches(0.8)
        )
        dep_box.fill.solid()
        dep_box.fill.fore_color.rgb = self.colors['ice_blue']
        dep_box.line.color.rgb = self.colors['secondary']

        dep_frame = dep_box.text_frame
        dep_frame.text = "📌 useEffect 依赖: [messages, sessionId, selectedTurnId]"
        p = dep_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # 性能优化提示
        perf_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.3), Inches(6.4),
            Inches(6.533), Inches(0.8)
        )
        perf_box.fill.solid()
        perf_box.fill.fore_color.rgb = self.colors['ice_blue']
        perf_box.line.color.rgb = self.colors['accent']

        perf_frame = perf_box.text_frame
        perf_frame.text = "⚡ 性能优化：字符串比较避免类型不匹配"
        p = perf_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['accent']

        self.add_slide_number(slide, "05")

    def add_interaction_slide(self):
        """用户交互流程"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self.add_header_bar(slide)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "04 用户交互流程 | User Interaction"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 三种触发方式
        triggers = [
            ("🖱️ 方式一", "点击 AI 消息\n（长消息按钮）", self.colors['secondary']),
            ("📝 方式二", "点击 AI 消息\n（短消息内容）", self.colors['accent']),
            ("👤 方式三", "点击用户消息", self.colors['success']),
        ]

        y_start = 1.6
        for i, (icon, text, color) in enumerate(triggers):
            x = 0.5 + i * 4.2
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y_start),
                Inches(4), Inches(1.8)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.colors['white']
            card.line.color.rgb = color
            card.line.width = Pt(3)

            # 文字
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y_start + 0.6),
                Inches(3.6), Inches(1)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = f"{icon}\n\n{text}"
            text_para = text_frame.paragraphs[0]
            text_para.alignment = PP_ALIGN.CENTER
            text_para.font.size = Pt(14)
            text_para.font.bold = True
            text_para.font.color.rgb = self.colors['dark']
            text_para.line_spacing = 1.3

        # 代码实现
        impl_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8), Inches(5.5), Inches(0.4)
        )
        it_frame = impl_title.text_frame
        it_frame.text = "💻 代码实现"
        p = it_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        impl_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(4.3),
            Inches(5.5), Inches(1.8)
        )
        impl_box.fill.solid()
        impl_box.fill.fore_color.rgb = self.colors['dark']
        impl_box.line.color.rgb = self.colors['secondary']
        impl_box.line.width = Pt(2)

        impl_frame = impl_box.text_frame
        impl_frame.word_wrap = True
        impl_frame.margin_left = Inches(0.2)

        impl_code = '''onClick={() => {
  if (m.conversation_turn_id) {
    setSelectedTurnId(
      m.conversation_turn_id
    );
  }
}}'''

        p = impl_frame.paragraphs[0]
        p.text = impl_code
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(100, 255, 100)
        p.font.name = 'Courier New'

        # 交互流程图
        flow_title = slide.shapes.add_textbox(
            Inches(6.3), Inches(3.8), Inches(6.533), Inches(0.4)
        )
        ft_frame = flow_title.text_frame
        ft_frame.text = "🔄 交互流程"
        p = ft_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        flow_steps = [
            ("用户点击消息", 6.3, 4.3, self.colors['secondary']),
            ("设置 selectedTurnId", 6.3, 4.95, self.colors['accent']),
            ("触发 useEffect", 6.3, 5.6, self.colors['success']),
            ("提取并显示工具", 6.3, 6.25, self.colors['primary']),
        ]

        for step, x, y, color in flow_steps:
            step_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(6.533), Inches(0.5)
            )
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = color
            step_box.line.fill.background()

            step_text = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.1),
                Inches(6.133), Inches(0.3)
            )
            st_frame = step_text.text_frame
            st_frame.text = step
            p = st_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

        # 位置信息
        loc_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.4),
            Inches(12.333), Inches(0.8)
        )
        loc_box.fill.solid()
        loc_box.fill.fore_color.rgb = self.colors['light_gray']
        loc_box.line.color.rgb = self.colors['gray']

        loc_frame = loc_box.text_frame
        loc_frame.text = "📍 代码位置：ChatInterface.tsx:2112-2171"
        p = loc_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']

        self.add_slide_number(slide, "06")

    def add_troubleshooting_slide(self):
        """问题诊断方案"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self.add_header_bar(slide)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "05 问题诊断方案 | Troubleshooting Guide"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        problems = [
            ("❌ 问题 1", "selectedTurnId 未设置",
             "点击消息无日志输出 → 检查 conversation_turn_id 字段",
             self.colors['warning']),

            ("❌ 问题 2", "useEffect 未触发",
             "状态变化无日志 → 确认依赖数组包含 selectedTurnId",
             self.colors['warning']),

            ("❌ 问题 3", "工具调用匹配失败",
             "检查匹配逻辑日志 → 验证类型转换和字符串比较",
             self.colors['warning']),

            ("❌ 问题 4", "后端数据缺失",
             "tool_calls 缺少 turn_id → 确保后端返回完整数据",
             self.colors['warning']),
        ]

        y_start = 1.5
        for i, (icon, title, desc, color) in enumerate(problems):
            y_pos = y_start + i * 1.2

            # 问题框
            prob_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y_pos),
                Inches(12.333), Inches(1.05)
            )
            prob_box.fill.solid()
            prob_box.fill.fore_color.rgb = self.colors['white']
            prob_box.line.color.rgb = color
            prob_box.line.width = Pt(2)

            # 图标和标题
            title_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y_pos + 0.2),
                Inches(3), Inches(0.3)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"{icon} {title}"
            p = title_frame.paragraphs[0]
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']

            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y_pos + 0.55),
                Inches(11.933), Inches(0.35)
            )
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_frame.text = desc
            p = desc_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['gray']

        # 调试工具提示
        tip_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.4),
            Inches(12.333), Inches(0.8)
        )
        tip_box.fill.solid()
        tip_box.fill.fore_color.rgb = self.colors['ice_blue']
        tip_box.line.color.rgb = self.colors['secondary']

        tip_frame = tip_box.text_frame
        tip_frame.text = "💡 调试工具：使用控制台日志快速定位问题，参考右侧检查清单"
        p = tip_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        self.add_slide_number(slide, "07")

    def add_code_reference_slide(self):
        """关键代码索引"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self.add_header_bar(slide)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "06 关键代码索引 | Code Reference"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 代码表格
        codes = [
            ("API 调用", "agentService.ts", "393"),
            ("数据转换", "agentService.ts", "414-448"),
            ("工具调用提取", "ChatInterface.tsx", "1368-1534"),
            ("点击消息设置", "ChatInterface.tsx", "2112-2171"),
            ("UI 渲染", "ChatInterface.tsx", "2518-2583"),
        ]

        # 表头
        headers = ["功能模块", "文件路径", "代码行号"]
        header_widths = [2.5, 6.5, 2.833]
        x_positions = [0.5, 3, 9.5]
        y_start = 1.5

        for i, (header, width, x) in enumerate(zip(headers, header_widths, x_positions)):
            header_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y_start),
                Inches(width), Inches(0.6)
            )
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = self.colors['primary']
            header_box.line.fill.background()

            header_frame = header_box.text_frame
            header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            header_frame.text = header
            p = header_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

        # 表格内容
        for i, (func, file, line) in enumerate(codes):
            y_pos = y_start + 0.6 + i * 0.75

            # 功能
            func_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(y_pos),
                Inches(2.5), Inches(0.7)
            )
            func_box.fill.solid()
            func_box.fill.fore_color.rgb = self.colors['light_gray']
            func_box.line.color.rgb = self.colors['gray']
            func_box.line.width = Pt(1)

            func_frame = func_box.text_frame
            func_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            func_frame.text = func
            p = func_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']

            # 文件
            file_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(3), Inches(y_pos),
                Inches(6.5), Inches(0.7)
            )
            file_box.fill.solid()
            file_box.fill.fore_color.rgb = self.colors['white']
            file_box.line.color.rgb = self.colors['gray']
            file_box.line.width = Pt(1)

            file_frame = file_box.text_frame
            file_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            file_frame.text = file
            p = file_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['dark']
            p.font.name = 'Courier New'

            # 行号
            line_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(9.5), Inches(y_pos),
                Inches(2.833), Inches(0.7)
            )
            line_box.fill.solid()
            line_box.fill.fore_color.rgb = self.colors['ice_blue']
            line_box.line.color.rgb = self.colors['gray']
            line_box.line.width = Pt(1)

            line_frame = line_box.text_frame
            line_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            line_frame.text = line
            p = line_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.colors['accent']

        # 底部提示
        tip_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.9), Inches(12.333), Inches(1.3)
        )
        tip_frame = tip_box.text_frame
        tip_frame.word_wrap = True
        tip_frame.text = "📌 提示：使用 Ctrl+F (或 Cmd+F) 在代码编辑器中快速定位这些行号"
        p = tip_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.color.rgb = self.colors['gray']
        p.font.italic = True

        self.add_slide_number(slide, "08")

    def add_summary_slide(self):
        """总结与建议"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 渐变背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['navy']
        bg.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "总结与建议"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['white']

        # 核心要点
        summary_items = [
            ("✅", "数据流完整", "API → 转换 → 存储 → 提取 → 渲染"),
            ("✅", "数据结构完善", "双重 ID 保障：消息级 + 工具级"),
            ("✅", "交互逻辑清晰", "点击 → 设置 → 提取 → 显示"),
            ("🔍", "问题定位明确", "重点检查提取匹配机制"),
            ("🛠️", "调试工具完备", "完善的日志系统支持"),
        ]

        y_start = 2.8
        for icon, title, desc in summary_items:
            y_pos = y_start + summary_items.index((icon, title, desc)) * 0.55

            # 图标
            icon_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_pos), Inches(0.4), Inches(0.4)
            )
            icon_frame = icon_box.text_frame
            icon_frame.text = icon
            p = icon_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['accent']

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(2.1), Inches(y_pos), Inches(2), Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            p = title_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']

            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(4.2), Inches(y_pos), Inches(8), Inches(0.4)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = desc
            p = desc_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(180, 200, 230)

        # 底部装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(5.6), Inches(3.333), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent']
        line.line.fill.background()

        # 联系信息
        contact_box = slide.shapes.add_textbox(
            Inches(1), Inches(6), Inches(11.333), Inches(1)
        )
        contact_frame = contact_box.text_frame
        contact_frame.word_wrap = True

        p = contact_frame.paragraphs[0]
        p.text = "基于完整代码分析与数据流梳理"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(150, 200, 255)

        p = contact_frame.add_paragraph()
        p.text = "Technical Documentation | Generated by AI Assistant"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(100, 150, 200)
        p.space_before = Pt(8)

        self.add_slide_number(slide, "09")

    def generate(self, output_path):
        """生成PPT"""
        self.add_title_slide()
        self.add_agenda_slide()
        self.add_data_flow_slide()
        self.add_data_acquisition_slide()
        self.add_extraction_slide()
        self.add_interaction_slide()
        self.add_troubleshooting_slide()
        self.add_code_reference_slide()
        self.add_summary_slide()

        self.prs.save(output_path)
        print(f"✅ 专业商业PPT已生成: {output_path}")
        print(f"📊 共 {len(self.prs.slides)} 页幻灯片")
        return output_path

if __name__ == "__main__":
    generator = ProfessionalPPTGenerator()
    output_file = "/Users/hehe/pycharm_projects/aigc/backend/work_dir/Tool_Calls_流程分析_专业版.pptx"
    generator.generate(output_file)
