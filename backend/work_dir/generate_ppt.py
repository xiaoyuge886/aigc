#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Tool Calls 流程分析 - 商业演示文稿生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

class ToolCallsPPTGenerator:
    """生成 Tool Calls 流程分析的商业PPT"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # 定义配色方案
        self.colors = {
            'primary': RGBColor(0, 82, 147),      # 深蓝色
            'secondary': RGBColor(0, 164, 239),    # 亮蓝色
            'accent': RGBColor(255, 127, 0),       # 橙色
            'success': RGBColor(40, 167, 69),      # 绿色
            'warning': RGBColor(255, 193, 7),      # 黄色
            'dark': RGBColor(52, 58, 64),          # 深灰色
            'light': RGBColor(248, 249, 250),      # 浅灰色
            'white': RGBColor(255, 255, 255),
            'gradient_start': RGBColor(0, 82, 147),
            'gradient_end': RGBColor(0, 164, 239)
        }

    def add_title_slide(self):
        """添加标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 背景渐变效果
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['primary']
        background.line.fill.background()

        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Tool Calls 流程分析"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['white']

        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2), Inches(11.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "完整数据流与实现机制解析"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(200, 220, 255)

        # 装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(4), Inches(3.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent']
        line.line.fill.background()

    def add_overview_slide(self):
        """添加数据流概览页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "📊 数据流概览"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 流程框
        flows = [
            "后端 API",
            "GET /api/v1/session/{sessionId}/conversation",
            "getConversationHistory()",
            "数据转换与格式化",
            "保留 conversation_turn_id 和 tool_calls",
            "ChatInterface 组件",
            "useEffect 监听与提取",
            "UI 渲染 Tool Calls"
        ]

        y_start = 1.8
        box_height = 0.6
        box_spacing = 0.15

        for i, flow_text in enumerate(flows):
            y_pos = y_start + i * (box_height + box_spacing)

            # 流程框
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(2), Inches(y_pos),
                Inches(9.333), Inches(box_height)
            )

            box.fill.solid()
            if i in [1, 2, 4]:
                box.fill.fore_color.rgb = self.colors['secondary']
            elif i in [5, 6]:
                box.fill.fore_color.rgb = self.colors['accent']
            else:
                box.fill.fore_color.rgb = self.colors['primary']

            box.line.color.rgb = self.colors['white']
            box.line.width = Pt(2)

            # 文本
            text_frame = box.text_frame
            text_frame.text = flow_text
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            para = text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = self.colors['white']

            # 添加箭头（除了最后一个）
            if i < len(flows) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.ISOSCELES_TRIANGLE,
                    Inches(6.5), Inches(y_pos + box_height + 0.05),
                    Inches(0.3), Inches(0.1)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['dark']
                arrow.line.fill.background()
                arrow.rotation = 180

    def add_data_flow_slide(self):
        """添加数据获取流程页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "1️⃣ 数据获取流程"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 左侧：API调用信息
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6), Inches(6), Inches(5.4)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        # 标题
        p = left_frame.paragraphs[0]
        p.text = "API 调用"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']

        # 内容
        items = [
            ("接口", "GET /api/v1/session/{sessionId}/conversation"),
            ("位置", "agentService.ts:390-458"),
            ("功能", "从后端获取对话历史"),
            ("触发时机", "• 初始化时\n• 切换会话时\n• 外部会话切换")
        ]

        for title, content in items:
            p = left_frame.add_paragraph()
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            p.space_before = Pt(12)

            p = left_frame.add_paragraph()
            p.text = content
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['dark']
            p.space_before = Pt(2)
            p.level = 1

        # 右侧：数据结构
        right_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), Inches(1.6),
            Inches(6), Inches(5.4)
        )
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = self.colors['light']
        right_box.line.color.rgb = self.colors['secondary']
        right_box.line.width = Pt(2)

        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        right_frame.margin_top = Inches(0.15)
        right_frame.margin_bottom = Inches(0.15)
        right_frame.margin_left = Inches(0.15)
        right_frame.margin_right = Inches(0.15)

        p = right_frame.paragraphs[0]
        p.text = "返回数据结构"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        code = '''{
  "session_id": "xxx",
  "messages": [
    {
      "id": "ai-3",
      "text": "...",
      "sender": "ai",
      "timestamp": "2025-12-30...",
      "conversation_turn_id": "98a93fe3...",
      "tool_calls": [
        {
          "tool_use_id": "call_xxx",
          "tool_name": "WebSearch",
          "tool_input": {...},
          "tool_output": null,
          "conversation_turn_id": "98a93fe3..."
        }
      ]
    }
  ]
}'''

        p = right_frame.add_paragraph()
        p.text = code
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors['dark']
        p.font.name = 'Courier New'
        p.space_before = Pt(8)

    def add_extraction_slide(self):
        """添加工具调用提取流程页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "2️⃣ 工具调用提取流程"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 位置信息
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5)
        )
        info_frame = info_box.text_frame
        info_frame.text = "位置：ChatInterface.tsx:1368-1534 | 依赖：[messages, sessionId, selectedTurnId]"
        info_para = info_frame.paragraphs[0]
        info_para.font.size = Pt(14)
        info_para.font.color.rgb = self.colors['dark']
        info_para.font.italic = True

        # 三步流程
        steps = [
            {
                "num": "1",
                "title": "确定目标轮次",
                "content": "• 检查是否有 selectedTurnId\n• 有：使用用户选择的轮次\n• 无：使用最新 AI 消息的轮次",
                "color": self.colors['secondary']
            },
            {
                "num": "2",
                "title": "提取工具调用",
                "content": "• 遍历所有消息的 tool_calls\n• 匹配 conversation_turn_id\n• 只提取目标轮次的工具调用",
                "color": self.colors['accent']
            },
            {
                "num": "3",
                "title": "更新状态",
                "content": "• setToolCalls(extractedToolCalls)\n• 触发 UI 重新渲染\n• 显示选中轮次的工具调用",
                "color": self.colors['success']
            }
        ]

        y_start = 2.3
        box_width = 3.8
        box_height = 2.8
        box_spacing = 0.3

        for i, step in enumerate(steps):
            x_pos = 0.5 + i * (box_width + box_spacing)

            # 步骤框
            step_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_start),
                Inches(box_width), Inches(box_height)
            )
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = self.colors['light']
            step_box.line.color.rgb = step['color']
            step_box.line.width = Pt(3)

            # 步骤编号
            num_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(y_start + 0.15),
                Inches(0.6), Inches(0.6)
            )
            num_frame = num_box.text_frame
            num_frame.text = step['num']
            num_para = num_frame.paragraphs[0]
            num_para.alignment = PP_ALIGN.CENTER
            num_para.font.size = Pt(32)
            num_para.font.bold = True
            num_para.font.color.rgb = step['color']

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.8), Inches(y_start + 0.2),
                Inches(2.8), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = step['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.color.rgb = step['color']

            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.15), Inches(y_start + 0.85),
                Inches(3.5), Inches(1.8)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = step['content']
            content_para = content_frame.paragraphs[0]
            content_para.font.size = Pt(13)
            content_para.font.color.rgb = self.colors['dark']
            content_para.line_spacing = 1.4

        # 底部关键逻辑
        logic_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(5.4),
            Inches(12.333), Inches(1.6)
        )
        logic_box.fill.solid()
        logic_box.fill.fore_color.rgb = self.colors['primary']
        logic_box.line.fill.background()

        logic_frame = logic_box.text_frame
        logic_frame.word_wrap = True
        logic_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        logic_frame.margin_left = Inches(0.3)
        logic_frame.margin_right = Inches(0.3)

        p = logic_frame.paragraphs[0]
        p.text = "🔑 核心匹配逻辑"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        logic_code = '''const toolTurnIdStr = String(toolCallTurnId || '');
const targetTurnIdStr = String(targetTurnId || '');
if (toolTurnIdStr === targetTurnIdStr) {
    // 匹配成功，添加到工具调用列表
}'''

        p = logic_frame.add_paragraph()
        p.text = logic_code
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.font.name = 'Courier New'
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

    def add_interaction_slide(self):
        """添加用户交互流程页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "3️⃣ 用户交互流程"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 位置信息
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5)
        )
        info_frame = info_box.text_frame
        info_frame.text = "位置：ChatInterface.tsx:2112-2171"
        info_para = info_frame.paragraphs[0]
        info_para.font.size = Pt(14)
        info_para.font.color.rgb = self.colors['dark']
        info_para.font.italic = True

        # 三种触发方式
        triggers = [
            {
                "icon": "🖱️",
                "title": "点击 AI 消息（长消息）",
                "desc": "点击长消息的按钮区域"
            },
            {
                "icon": "📝",
                "title": "点击 AI 消息（短消息）",
                "desc": "直接点击短消息的内容区域"
            },
            {
                "icon": "👤",
                "title": "点击用户消息",
                "desc": "点击用户发送的消息"
            }
        ]

        y_start = 2.3
        box_width = 3.8
        box_height = 1.8
        box_spacing = 0.3

        for i, trigger in enumerate(triggers):
            x_pos = 0.5 + i * (box_width + box_spacing)

            # 触发框
            trigger_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_start),
                Inches(box_width), Inches(box_height)
            )
            trigger_box.fill.solid()
            trigger_box.fill.fore_color.rgb = self.colors['light']
            trigger_box.line.color.rgb = self.colors['secondary']
            trigger_box.line.width = Pt(2)

            # 图标和标题
            header_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y_start + 0.2),
                Inches(3.4), Inches(0.6)
            )
            header_frame = header_box.text_frame
            header_frame.text = f"{trigger['icon']} {trigger['title']}"
            header_para = header_frame.paragraphs[0]
            header_para.font.size = Pt(16)
            header_para.font.bold = True
            header_para.font.color.rgb = self.colors['primary']

            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y_start + 0.9),
                Inches(3.4), Inches(0.7)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = trigger['desc']
            desc_frame.word_wrap = True
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(13)
            desc_para.font.color.rgb = self.colors['dark']

        # 设置逻辑
        logic_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(4.4),
            Inches(12.333), Inches(2.6)
        )
        logic_box.fill.solid()
        logic_box.fill.fore_color.rgb = RGBColor(232, 244, 255)
        logic_box.line.color.rgb = self.colors['secondary']
        logic_box.line.width = Pt(2)

        logic_title = slide.shapes.add_textbox(
            Inches(0.7), Inches(4.6),
            Inches(12), Inches(0.5)
        )
        logic_title_frame = logic_title.text_frame
        logic_title_frame.text = "⚙️ 设置逻辑"
        p = logic_title_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        logic_code = '''onClick={() => {
  if (m.conversation_turn_id) {
    setSelectedTurnId(m.conversation_turn_id);
  }
}'''

        logic_content = slide.shapes.add_textbox(
            Inches(0.7), Inches(5.2),
            Inches(12), Inches(1.6)
        )
        logic_content_frame = logic_content.text_frame
        logic_content_frame.text = logic_code
        logic_content_frame.word_wrap = True
        p = logic_content_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['dark']
        p.font.name = 'Courier New'

    def add_troubleshooting_slide(self):
        """添加问题诊断页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "4️⃣ 常见问题诊断"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        problems = [
            {
                "title": "❌ selectedTurnId 未正确设置",
                "check": "检查：点击消息时是否有日志输出",
                "solve": "解决：确保消息有 conversation_turn_id 字段",
                "color": self.colors['warning']
            },
            {
                "title": "❌ useEffect 未触发",
                "check": "检查：selectedTurnId 变化时是否有日志",
                "solve": "解决：确保依赖数组包含 selectedTurnId",
                "color": self.colors['warning']
            },
            {
                "title": "❌ 工具调用匹配失败",
                "check": "检查：工具调用轮次匹配检查日志",
                "solve": "解决：检查类型转换和字符串比较",
                "color": self.colors['warning']
            },
            {
                "title": "❌ tool_calls 缺少 conversation_turn_id",
                "check": "检查：getConversationHistory 消息日志",
                "solve": "解决：确保后端返回包含 conversation_turn_id",
                "color": self.colors['warning']
            }
        ]

        y_start = 1.6
        box_height = 1.3
        box_spacing = 0.15

        for i, problem in enumerate(problems):
            y_pos = y_start + i * (box_height + box_spacing)

            # 问题框
            problem_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y_pos),
                Inches(12.333), Inches(box_height)
            )
            problem_box.fill.solid()
            problem_box.fill.fore_color.rgb = self.colors['light']
            problem_box.line.color.rgb = problem['color']
            problem_box.line.width = Pt(2)

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y_pos + 0.15),
                Inches(4), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = problem['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(16)
            title_para.font.bold = True
            title_para.font.color.rgb = problem['color']

            # 检查项
            check_box = slide.shapes.add_textbox(
                Inches(4.5), Inches(y_pos + 0.15),
                Inches(4), Inches(0.9)
            )
            check_frame = check_box.text_frame
            check_frame.word_wrap = True
            check_frame.text = f"🔍 {problem['check']}"
            check_para = check_frame.paragraphs[0]
            check_para.font.size = Pt(13)
            check_para.font.color.rgb = self.colors['dark']

            # 解决方案
            solve_box = slide.shapes.add_textbox(
                Inches(8.8), Inches(y_pos + 0.15),
                Inches(3.8), Inches(0.9)
            )
            solve_frame = solve_box.text_frame
            solve_frame.word_wrap = True
            solve_frame.text = f"✅ {problem['solve']}"
            solve_para = solve_frame.paragraphs[0]
            solve_para.font.size = Pt(13)
            solve_para.font.color.rgb = self.colors['success']

    def add_checklist_slide(self):
        """添加调试检查清单页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "5️⃣ 调试检查清单"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        checks = [
            {
                "step": "1",
                "title": "检查数据获取",
                "logs": "[getConversationHistory] Fetching history...\n[getConversationHistory] Response data: {...}\n[getConversationHistory] 消息 #X: {conversation_turn_id: \"xxx\"}"
            },
            {
                "step": "2",
                "title": "检查消息渲染",
                "logs": "[ChatInterface] Rendering message: {\n  conversation_turn_id: \"xxx\",\n  has_tool_calls: true,\n  tool_calls_count: X\n}"
            },
            {
                "step": "3",
                "title": "检查点击事件",
                "logs": "🖱️ 点击短消息: {conversation_turn_id: \"xxx\"}\n✅ 已设置 selectedTurnId: {conversation_turn_id: \"xxx\"}"
            },
            {
                "step": "4",
                "title": "检查状态变化",
                "logs": "🎯 selectedTurnId 状态变化: {selectedTurnId: \"xxx\"}\n🔧 [useEffect] 开始提取工具调用"
            },
            {
                "step": "5",
                "title": "检查工具调用匹配",
                "logs": "🔍 工具调用轮次匹配检查: {\n  tool_turn_id: \"xxx\",\n  target_turn_id: \"xxx\",\n  isMatch: true/false\n}"
            },
            {
                "step": "6",
                "title": "检查最终状态",
                "logs": "🔧 工具调用状态更新: {total: X, selectedTurnId: \"xxx\"}\n🔧 工具列表渲染状态: {toolCallsCount: X}"
            }
        ]

        # 左列
        left_col = checks[:3]
        for i, check in enumerate(left_col):
            y_pos = 1.6 + i * 1.75

            # 步骤框
            step_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y_pos),
                Inches(6), Inches(1.6)
            )
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = self.colors['light']
            step_box.line.color.rgb = self.colors['secondary']
            step_box.line.width = Pt(2)

            # 步骤编号
            num_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y_pos + 0.1),
                Inches(0.5), Inches(0.5)
            )
            num_frame = num_box.text_frame
            num_frame.text = check['step']
            num_para = num_frame.paragraphs[0]
            num_para.font.size = Pt(24)
            num_para.font.bold = True
            num_para.font.color.rgb = self.colors['accent']

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(1.3), Inches(y_pos + 0.1),
                Inches(5), Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = check['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(16)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['primary']

            # 日志内容
            logs_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y_pos + 0.55),
                Inches(5.6), Inches(0.95)
            )
            logs_frame = logs_box.text_frame
            logs_frame.word_wrap = True
            logs_frame.text = check['logs']
            logs_para = logs_frame.paragraphs[0]
            logs_para.font.size = Pt(9)
            logs_para.font.color.rgb = self.colors['dark']
            logs_para.font.name = 'Courier New'

        # 右列
        right_col = checks[3:]
        for i, check in enumerate(right_col):
            y_pos = 1.6 + i * 1.75

            # 步骤框
            step_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.8), Inches(y_pos),
                Inches(6), Inches(1.6)
            )
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = self.colors['light']
            step_box.line.color.rgb = self.colors['secondary']
            step_box.line.width = Pt(2)

            # 步骤编号
            num_box = slide.shapes.add_textbox(
                Inches(7), Inches(y_pos + 0.1),
                Inches(0.5), Inches(0.5)
            )
            num_frame = num_box.text_frame
            num_frame.text = check['step']
            num_para = num_frame.paragraphs[0]
            num_para.font.size = Pt(24)
            num_para.font.bold = True
            num_para.font.color.rgb = self.colors['accent']

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(7.6), Inches(y_pos + 0.1),
                Inches(5), Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = check['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(16)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['primary']

            # 日志内容
            logs_box = slide.shapes.add_textbox(
                Inches(7), Inches(y_pos + 0.55),
                Inches(5.6), Inches(0.95)
            )
            logs_frame = logs_box.text_frame
            logs_frame.word_wrap = True
            logs_frame.text = check['logs']
            logs_para = logs_frame.paragraphs[0]
            logs_para.font.size = Pt(9)
            logs_para.font.color.rgb = self.colors['dark']
            logs_para.font.name = 'Courier New'

    def add_code_summary_slide(self):
        """添加关键代码位置页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "📁 关键代码位置"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 代码位置表格
        code_locations = [
            {"功能": "API 调用", "文件": "agentService.ts", "行号": "393"},
            {"功能": "数据转换", "文件": "agentService.ts", "行号": "414-448"},
            {"功能": "工具调用提取", "文件": "ChatInterface.tsx", "行号": "1368-1534"},
            {"功能": "点击消息设置轮次", "文件": "ChatInterface.tsx", "行号": "2112-2171"},
            {"功能": "UI 渲染", "文件": "ChatInterface.tsx", "行号": "2518-2583"}
        ]

        y_start = 1.6
        row_height = 0.9

        # 表头
        headers = ["功能", "文件", "行号"]
        header_widths = [3, 5.5, 3.5]
        x_positions = [0.5, 3.5, 9]

        for i, header in enumerate(headers):
            header_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_positions[i]), Inches(y_start),
                Inches(header_widths[i]), Inches(row_height)
            )
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = self.colors['primary']
            header_box.line.fill.background()

            header_frame = header_box.text_frame
            header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            header_frame.text = header
            header_para = header_frame.paragraphs[0]
            header_para.alignment = PP_ALIGN.CENTER
            header_para.font.size = Pt(18)
            header_para.font.bold = True
            header_para.font.color.rgb = self.colors['white']

        # 表格内容
        for i, location in enumerate(code_locations):
            y_pos = y_start + row_height + i * row_height

            # 功能
            func_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(y_pos),
                Inches(3), Inches(row_height - 0.05)
            )
            func_box.fill.solid()
            func_box.fill.fore_color.rgb = self.colors['light']
            func_box.line.color.rgb = RGBColor(200, 200, 200)
            func_box.line.width = Pt(1)

            func_frame = func_box.text_frame
            func_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            func_frame.text = location["功能"]
            func_para = func_frame.paragraphs[0]
            func_para.alignment = PP_ALIGN.CENTER
            func_para.font.size = Pt(14)
            func_para.font.bold = True
            func_para.font.color.rgb = self.colors['dark']

            # 文件
            file_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(3.5), Inches(y_pos),
                Inches(5.5), Inches(row_height - 0.05)
            )
            file_box.fill.solid()
            file_box.fill.fore_color.rgb = self.colors['white']
            file_box.line.color.rgb = RGBColor(200, 200, 200)
            file_box.line.width = Pt(1)

            file_frame = file_box.text_frame
            file_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            file_frame.text = location["文件"]
            file_para = file_frame.paragraphs[0]
            file_para.alignment = PP_ALIGN.CENTER
            file_para.font.size = Pt(13)
            file_para.font.color.rgb = self.colors['dark']
            file_para.font.name = 'Courier New'

            # 行号
            line_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(9), Inches(y_pos),
                Inches(3.833), Inches(row_height - 0.05)
            )
            line_box.fill.solid()
            line_box.fill.fore_color.rgb = self.colors['white']
            line_box.line.color.rgb = RGBColor(200, 200, 200)
            line_box.line.width = Pt(1)

            line_frame = line_box.text_frame
            line_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            line_frame.text = location["行号"]
            line_para = line_frame.paragraphs[0]
            line_para.alignment = PP_ALIGN.CENTER
            line_para.font.size = Pt(14)
            line_para.font.bold = True
            line_para.font.color.rgb = self.colors['accent']
            line_para.font.name = 'Courier New'

    def add_data_validation_slide(self):
        """添加数据验证页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "✅ 数据验证"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        # 验证结果
        result_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(12.333), Inches(1.5)
        )
        result_box.fill.solid()
        result_box.fill.fore_color.rgb = self.colors['success']
        result_box.line.fill.background()

        result_frame = result_box.text_frame
        result_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        result_frame.text = "✅ 结论：数据完整，问题应在提取或匹配逻辑"
        result_para = result_frame.paragraphs[0]
        result_para.alignment = PP_ALIGN.CENTER
        result_para.font.size = Pt(24)
        result_para.font.bold = True
        result_para.font.color.rgb = self.colors['white']

        # 消息验证
        msg_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.4), Inches(6), Inches(3.6)
        )
        msg_frame = msg_box.text_frame
        msg_frame.word_wrap = True

        p = msg_frame.paragraphs[0]
        p.text = "📨 有 conversation_turn_id 的消息"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.space_after = Pt(12)

        messages = [
            "user-3: 98a93fe3e096408d",
            "ai-3: 98a93fe3e096408d",
            "user-4: 4c1ca2cc75254fa4",
            "ai-4: 4c1ca2cc75254fa4",
            "user-5: 69eb8567759e424c",
            "ai-5: 69eb8567759e424c"
        ]

        for msg in messages:
            p = msg_frame.add_paragraph()
            p.text = f"✓ {msg}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['dark']
            p.space_before = Pt(4)

        # 工具调用验证
        tool_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), Inches(3.4),
            Inches(6), Inches(3.6)
        )
        tool_box.fill.solid()
        tool_box.fill.fore_color.rgb = self.colors['light']
        tool_box.line.color.rgb = self.colors['secondary']
        tool_box.line.width = Pt(2)

        tool_frame = tool_box.text_frame
        tool_frame.word_wrap = True
        tool_frame.margin_top = Inches(0.2)
        tool_frame.margin_left = Inches(0.2)
        tool_frame.margin_right = Inches(0.2)

        p = tool_frame.paragraphs[0]
        p.text = "🔧 tool_calls 中也有 conversation_turn_id"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.space_after = Pt(12)

        code = '''{
  "tool_use_id": "call_xxx",
  "tool_name": "WebSearch",
  "tool_input": {...},
  "tool_output": null,
  "conversation_turn_id": "98a93fe3..."
}'''

        p = tool_frame.add_paragraph()
        p.text = code
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['dark']
        p.font.name = 'Courier New'

    def add_next_steps_slide(self):
        """添加下一步调试页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "🚀 下一步调试步骤"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']

        steps = [
            {
                "num": "1",
                "title": "刷新页面",
                "desc": "查看控制台是否有 [getConversationHistory] 相关日志",
                "icon": "🔄"
            },
            {
                "num": "2",
                "title": "点击有 conversation_turn_id 的消息",
                "desc": "例如 user-3 或 ai-3",
                "icon": "🖱️"
            },
            {
                "num": "3",
                "title": "查看控制台输出",
                "desc": "确认点击事件和状态变化日志",
                "icon": "🔍"
            },
            {
                "num": "4",
                "title": "切换到 Tool Calls 标签页",
                "desc": "查看是否显示对应轮次的工具调用",
                "icon": "📋"
            }
        ]

        y_start = 1.7
        box_height = 1.3
        box_spacing = 0.2

        for i, step in enumerate(steps):
            y_pos = y_start + i * (box_height + box_spacing)

            # 步骤框
            step_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y_pos),
                Inches(12.333), Inches(box_height)
            )
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = self.colors['light']
            step_box.line.color.rgb = self.colors['secondary']
            step_box.line.width = Pt(2)

            # 编号
            num_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y_pos + 0.2),
                Inches(0.6), Inches(0.6)
            )
            num_frame = num_box.text_frame
            num_frame.text = step['num']
            num_para = num_frame.paragraphs[0]
            num_para.alignment = PP_ALIGN.CENTER
            num_para.font.size = Pt(32)
            num_para.font.bold = True
            num_para.font.color.rgb = self.colors['accent']

            # 图标和标题
            title_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_pos + 0.2),
                Inches(4.5), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"{step['icon']} {step['title']}"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['primary']

            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_pos + 0.7),
                Inches(11), Inches(0.5)
            )
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_frame.text = step['desc']
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(14)
            desc_para.font.color.rgb = self.colors['dark']

        # 底部提示
        tip_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.4),
            Inches(12.333), Inches(0.7)
        )
        tip_box.fill.solid()
        tip_box.fill.fore_color.rgb = self.colors['warning']
        tip_box.line.fill.background()

        tip_frame = tip_box.text_frame
        tip_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tip_frame.text = "💡 如果还是没有输出，请提供控制台的完整日志"
        tip_para = tip_frame.paragraphs[0]
        tip_para.alignment = PP_ALIGN.CENTER
        tip_para.font.size = Pt(16)
        tip_para.font.bold = True
        tip_para.font.color.rgb = self.colors['dark']

    def add_summary_slide(self):
        """添加总结页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 背景渐变
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['primary']
        background.line.fill.background()

        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.333), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = "总结"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['white']

        # 总结要点
        summary_points = [
            "✅ 数据流完整：API → 转换 → 存储 → 提取 → 渲染",
            "✅ 数据结构完整：conversation_turn_id 和 tool_calls 都有值",
            "✅ 交互逻辑清晰：点击消息 → 设置轮次 → 提取工具 → UI渲染",
            "🔍 问题定位：重点检查提取逻辑和匹配机制",
            "🛠️ 调试工具：完善的日志系统支持快速定位问题"
        ]

        y_start = 3.5
        for i, point in enumerate(summary_points):
            y_pos = y_start + i * 0.5

            point_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_pos),
                Inches(10.333), Inches(0.5)
            )
            point_frame = point_box.text_frame
            point_frame.text = point
            point_para = point_frame.paragraphs[0]
            point_para.font.size = Pt(20)
            point_para.font.color.rgb = RGBColor(200, 220, 255)

        # 底部装饰
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(5.8), Inches(3.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['accent']
        line.line.fill.background()

        # 联系方式
        contact_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.2), Inches(11.333), Inches(0.8)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = "基于完整的代码分析和数据流梳理"
        contact_para = contact_frame.paragraphs[0]
        contact_para.alignment = PP_ALIGN.CENTER
        contact_para.font.size = Pt(16)
        contact_para.font.color.rgb = RGBColor(150, 200, 255)
        contact_para.font.italic = True

    def generate(self, output_path):
        """生成PPT"""
        # 添加所有幻灯片
        self.add_title_slide()
        self.add_overview_slide()
        self.add_data_flow_slide()
        self.add_extraction_slide()
        self.add_interaction_slide()
        self.add_troubleshooting_slide()
        self.add_checklist_slide()
        self.add_code_summary_slide()
        self.add_data_validation_slide()
        self.add_next_steps_slide()
        self.add_summary_slide()

        # 保存文件
        self.prs.save(output_path)
        print(f"✅ PPT 已生成: {output_path}")
        return output_path

if __name__ == "__main__":
    # 生成PPT
    generator = ToolCallsPPTGenerator()
    output_file = "/Users/hehe/pycharm_projects/aigc/backend/work_dir/Tool_Calls_流程分析.pptx"
    generator.generate(output_file)
